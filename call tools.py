import os
from pptx import Presentation, util
from pptx.util import Inches, Pt

#图片位置设置
leftimg = util.Cm(10.61)
topimg = util.Cm(1.39)
widthimg = util.Cm(11)
heightimg = util.Cm(11)
#文本框位置设置
lefttext = util.Cm(11.30)
toptext = util.Cm(12.94)
widthtext = util.Cm(9.39)
heighttext = util.Cm(4)

#path_imgs = input()
path_imgs = './pictures'

# 实例化一个ppt演示文稿对象
prs = Presentation()
# 调整页面大小
prs.slide_width = util.Cm(32)
prs.slide_height = util.Cm(19)
# 实例化空白模板
blank_slide_layout = prs.slide_layouts[6]
#运行主代码部分
for files in os.listdir(path_imgs):
    img = files[0:-4]
    print(img)
    where = path_imgs + '/' + files
    print(where)
    slide = prs.slides.add_slide(blank_slide_layout)
    #图片
    pic = slide.shapes.add_picture(where, leftimg, topimg, widthimg, heightimg)
    #文本
    textbox=slide.shapes.add_textbox(lefttext,toptext,widthtext,heighttext)
    #设置字体
    ft = textbox.text_frame
    p = ft.add_paragraph()
    p.font.bold = True    # 加粗
    p.font.size = Pt(60)    # 大小
    p.text= img

save = input("保存为：") + '.pptx'
prs.save(save)